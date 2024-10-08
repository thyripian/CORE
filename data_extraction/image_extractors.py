import logging
import os
import shutil
from io import BytesIO
from zipfile import ZipFile

import fitz  # PyMuPDF
from docx import Document
from PIL import Image
from pptx import Presentation

logger = logging.getLogger(__name__)


# Function to compress an image using the WebP format, which can be set to lossless compression.
def compress_image(image_bytes_io, lossless=True):
    """
    Compress an image using the WebP format and return the compressed image as a BytesIO object.

    Parameters:
    image_bytes_io (BytesIO): A BytesIO object containing image data.
    lossless (bool): Indicates if the compression should be lossless. Default is True.

    Returns:
    BytesIO: A BytesIO object containing the compressed image data.
    """
    # Open the image from the BytesIO object
    image = Image.open(image_bytes_io)

    # Create a new BytesIO object to store the compressed image
    compressed_image_io = BytesIO()

    # Save the image in the WebP format with the specified compression type
    image.save(compressed_image_io, format="WebP", lossless=lossless)

    # Rewind the BytesIO object to the beginning for further reading
    compressed_image_io.seek(0)

    return compressed_image_io  # Return the compressed image BytesIO object


# Function to extract and compress images from a Word document, iterating through document relations.
def extract_images_from_word(file_path):
    """
    Extract images from a Word document and return them in a list of BytesIO objects.

    Parameters:
    file_path (str): The file path of the Word document.

    Returns:
    list: A list of BytesIO objects containing image data.
    """
    doc = Document(file_path)  # Open the Word document
    images = []  # Initialize an empty list to store images

    # Iterate over the relationships in the document
    for rel_id, rel in doc.part.rels.items():
        # Check if the relationship type is an image
        if "image" in rel.reltype:
            image_part = rel.target_part  # Get the part that contains the image
            image_bytes = image_part._blob  # Extract the image bytes
            image_bytes_io = BytesIO(
                image_bytes
            )  # Create a BytesIO object from the image bytes
            compressed_image = compress_image(image_bytes_io)  # Compress the image

            # Calculate the sizes of the compressed and original images
            compressed_size = len(compressed_image.getbuffer())
            original_size = len(image_bytes_io.getbuffer())

            # Add the smaller of the two images to the list
            if compressed_size > original_size:
                images.append(image_bytes_io)
            else:
                images.append(compressed_image)

    return images  # Return the list of image BytesIO objects


# Function to extract and compress images from an Excel file, handling the unique structure of Excel files.
def extract_images_from_excel(file_path):
    """
    Extract and compress images from an Excel file.

    This function first unzips the Excel file since it's a package of XML files and media.
    It then searches for image files in the 'xl/media' directory, compresses them using the
    compress_image function, and returns them as a list of BytesIO objects.

    Parameters:
    file_path (str): The path to the Excel file.

    Returns:
    list: A list of BytesIO objects containing the compressed image data.
    """
    temp_dir = "temp_xlsx"  # Define a temporary directory name
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)  # Create the directory if it doesn't exist

    images = []  # Initialize an empty list to store images

    # Use the ZipFile module to open and extract the Excel file
    with ZipFile(file_path, "r") as zip_ref:
        zip_ref.extractall(
            temp_dir
        )  # Extract all contents into the temporary directory

        # The images are usually contained within the 'xl/media' folder inside the unzipped directory
        media_path = os.path.join(temp_dir, "xl", "media")

        # Check if the media path exists
        if os.path.exists(media_path):
            # Iterate over each file in the media directory
            for filename in os.listdir(media_path):
                # Check if the file is an image
                if filename.endswith((".png", ".jpg", ".jpeg", ".gif")):
                    # Read the image file and create a BytesIO object
                    with open(os.path.join(media_path, filename), "rb") as f:
                        image_bytes_io = BytesIO(f.read())
                        compressed_image = compress_image(image_bytes_io)

                        # Calculate the size of the compressed and original image
                        compressed_size = len(compressed_image.getbuffer())
                        original_size = len(image_bytes_io.getbuffer())

                        # Choose the smaller image to append to the list
                        if compressed_size > original_size:
                            images.append(image_bytes_io)
                        else:
                            images.append(compressed_image)

    # Clean up by removing the temporary directory
    shutil.rmtree(temp_dir)

    return images  # Return the list of image BytesIO objects


# Function to extract and compress images from a PowerPoint file by iterating through slides and shapes.
def extract_images_from_pptx(file_path):
    """
    Extract and compress images from a PowerPoint (.pptx) file.

    This function iterates through the slides and shapes in a PowerPoint file to find images,
    compresses them using the compress_image function, and returns them as a list of BytesIO objects.

    Parameters:
    file_path (str): The path to the PowerPoint file.

    Returns:
    list: A list of BytesIO objects containing the compressed image data.
    """
    prs = Presentation(file_path)  # Open the PowerPoint file using python-pptx
    images = []  # Initialize an empty list to store images

    # Iterate over each slide in the presentation
    for slide in prs.slides:
        # Iterate over each shape in the slide
        for shape in slide.shapes:
            if (
                shape.shape_type == 13
            ):  # Check if the shape is an image (13 is the code for image shape)
                image_bytes = shape.image.blob  # Get the image blob from the shape
                image_bytes_io = BytesIO(
                    image_bytes
                )  # Create a BytesIO object from the blob
                compressed_image = compress_image(image_bytes_io)  # Compress the image

                # Calculate the size of the compressed and original image
                compressed_size = len(compressed_image.getbuffer())
                original_size = len(image_bytes_io.getbuffer())

                # Choose the smaller image to append to the list
                if compressed_size > original_size:
                    logger.info("\t\tExtracting as original image.")
                    images.append(image_bytes_io)
                else:
                    logger.info("\t\tExtracting as WebP compressed image.")
                    images.append(compressed_image)

    return images  # Return the list of image BytesIO objects


# Function to extract and compress images from a PDF file using the PyMuPDF library for direct image extraction.
def extract_images_from_pdf(file_path):
    """
    Extract and compress images from a PDF file using the PyMuPDF library.

    This function iterates over each page in the PDF, finds images, and compresses them.
    The compressed images are returned as a list of BytesIO objects.

    Parameters:
    file_path (str): The path to the PDF file.

    Returns:
    list: A list of BytesIO objects containing the compressed image data.
    """
    pdf_file = fitz.open(file_path)  # Open the PDF file using PyMuPDF
    images = []  # Initialize an empty list to store images

    #     logging.info(f"\tPDF has > {len(pdf_file)} < pages.") # Moved to READ_PDF FUNCTION above

    # Iterate over each page in the PDF
    for page_number in range(len(pdf_file)):
        page = pdf_file[page_number]  # Get the current page
        img_list = page.get_images(full=True)  # Get a list of images on the page

        # Iterate over each image in the list
        for img_index in img_list:
            xref = img_index[0]  # Get the xref of the image
            image = pdf_file.extract_image(xref)  # Extract the image using its xref
            image_bytes = image["image"]  # Get the image bytes
            image_bytes_io = BytesIO(
                image_bytes
            )  # Create a BytesIO object from the image bytes
            compressed_image = compress_image(image_bytes_io)  # Compress the image

            # Calculate the size of the compressed and original image
            compressed_size = len(compressed_image.getbuffer())
            original_size = len(image_bytes_io.getbuffer())

            # Choose the smaller image to append to the list
            if compressed_size > original_size:
                images.append(image_bytes_io)
            else:
                images.append(compressed_image)

    return images  # Return the list of image BytesIO objects
