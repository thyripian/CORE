import hashlib
import logging

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)


# read_word_file
def read_word_file(file_path):
    """
    Read a Word (.docx) file and extract all the text from it.

    Parameters:
    file_path (str): The path to the .docx file to read.

    Returns:
    str: The extracted text from the .docx file.
    """
    doc = Document(file_path)
    full_text = [
        paragraph.text for paragraph in doc.paragraphs
    ]  # Extract text from each paragraph
    return " ".join(full_text)


# read_excel_file
def read_excel_file(file_path):
    """
    Read an Excel (.xlsx) file and extract all text from the first worksheet.

    Parameters:
    file_path (str): The path to the .xlsx file to read.

    Returns:
    str: The extracted text from the .xlsx file.
    """
    wb = load_workbook(file_path)
    ws = wb.active
    full_text = [
        str(cell.value)
        for row in ws.iter_rows()
        for cell in row
        if cell.value is not None
    ]  # Extract text from each cell
    return " ".join(full_text)


# read_pptx_file
def read_pptx_file(file_path):
    """
    Read a PowerPoint (.pptx) file and extract all text from it.

    Parameters:
    file_path (str): The path to the .pptx file to read.

    Returns:
    str: The extracted text from the .pptx file.
    """
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text.extend(
                        run.text for run in paragraph.runs if run.text
                    )  # Extract text from each run
    return " ".join(full_text)


# Updated read_pdf_file
def read_pdf_file(file_path):
    """
    Read a PDF file and extract all unique text from it using the pdfplumber package.

    Parameters:
    file_path (str): The path to the PDF file to read.

    Returns:
    str: The extracted unique text from the PDF file.
    """
    full_text = []
    unique_pages = set()  # Use a set to track unique page hashes

    with pdfplumber.open(file_path) as pdf:
        logger.info(f"\t\tPDF has > {len(pdf.pages)} < pages.")

        for page in pdf.pages:
            text = page.extract_text()
            if text:
                # Create a hash of the page text
                page_hash = hashlib.md5(text.encode("utf-8")).hexdigest()

                # Only add the text if it's not already in the unique_pages set
                if page_hash not in unique_pages:
                    unique_pages.add(page_hash)
                    full_text.append(text)

    return " ".join(full_text)
