import os

import fitz  # PyMuPDF
import pdfplumber
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from file_handling import (read_excel_file, read_pdf_file, read_pptx_file,
                           read_word_file)


@pytest.fixture(scope="module")
def setup_files():
    files = {}

    # Create synthetic valid data
    files["valid_word_file"] = "valid_test.docx"
    files["valid_excel_file"] = "valid_test.xlsx"
    files["valid_pptx_file"] = "valid_test.pptx"
    files["valid_pdf_file"] = "valid_test.pdf"

    # Create a valid Word file with text
    doc = Document()
    doc.add_paragraph("This is a test document.")
    doc.save(files["valid_word_file"])

    # Create a valid Excel file with text
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "This is a test"
    wb.save(files["valid_excel_file"])

    # Create a valid PowerPoint file with text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "This is a test slide"
    prs.save(files["valid_pptx_file"])

    # Create a valid PDF file with text
    pdf_document = fitz.open()
    page = pdf_document.new_page()
    page.insert_text((100, 100), "This is a test page")
    pdf_document.save(files["valid_pdf_file"])

    # Create synthetic invalid data (empty files)
    files["empty_word_file"] = "empty_test.docx"
    files["empty_excel_file"] = "empty_test.xlsx"
    files["empty_pptx_file"] = "empty_test.pptx"
    files["empty_pdf_file"] = "empty_test.pdf"

    open(files["empty_word_file"], "w").close()
    open(files["empty_excel_file"], "w").close()
    open(files["empty_pptx_file"], "w").close()
    open(files["empty_pdf_file"], "w").close()

    # Create synthetic invalid data (corrupted files)
    files["corrupted_word_file"] = "corrupted_test.docx"
    files["corrupted_excel_file"] = "corrupted_test.xlsx"
    files["corrupted_pptx_file"] = "corrupted_test.pptx"
    files["corrupted_pdf_file"] = "corrupted_test.pdf"

    with open(files["corrupted_word_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_excel_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_pptx_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_pdf_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")

    yield files

    # Cleanup
    for file in files.values():
        if os.path.exists(file):
            os.remove(file)


@pytest.mark.parametrize(
    "file_key,read_function",
    [
        ("valid_word_file", read_word_file),
        ("valid_excel_file", read_excel_file),
        ("valid_pptx_file", read_pptx_file),
        ("valid_pdf_file", read_pdf_file),
    ],
)
def test_read_files_valid(setup_files, file_key, read_function):
    text = read_function(setup_files[file_key])
    assert len(text) > 0


@pytest.mark.parametrize(
    "file_key,read_function",
    [
        ("empty_word_file", read_word_file),
        ("empty_excel_file", read_excel_file),
        ("empty_pptx_file", read_pptx_file),
        ("empty_pdf_file", read_pdf_file),
    ],
)
def test_read_files_empty(setup_files, file_key, read_function):
    with pytest.raises(Exception):
        read_function(setup_files[file_key])


@pytest.mark.parametrize(
    "file_key,read_function",
    [
        ("corrupted_word_file", read_word_file),
        ("corrupted_excel_file", read_excel_file),
        ("corrupted_pptx_file", read_pptx_file),
        ("corrupted_pdf_file", read_pdf_file),
    ],
)
def test_read_files_corrupted(setup_files, file_key, read_function):
    with pytest.raises(Exception):
        read_function(setup_files[file_key])
