import os
import pytest
import pdfminer
from unittest.mock import patch, MagicMock
from zipfile import BadZipFile
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
import fitz  # PyMuPDF
from docx.opc.exceptions import PackageNotFoundError
from data_processing.data_processing_manager import DataProcessingManager
from data_processing.info_processing import extract_info
from data_processing.text_processing import (
    unique_subjects,
    extract_topics,
    normalize_mgrs,
    extract_classification_and_caveats,
    extract_keywords,
    filter_text
)
from utilities.keyword_loading import Keywords

@pytest.fixture(scope="module")
def setup_files():
    files = {}
    
    # Create synthetic valid data
    files['valid_word_file'] = 'valid_test.docx'
    files['valid_excel_file'] = 'valid_test.xlsx'
    files['valid_pptx_file'] = 'valid_test.pptx'
    files['valid_pdf_file'] = 'valid_test.pdf'

    # Create a valid Word file with text
    doc = Document()
    doc.add_paragraph("This is a test document.")
    doc.save(files['valid_word_file'])

    # Create a valid Excel file with text
    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'This is a test'
    wb.save(files['valid_excel_file'])

    # Create a valid PowerPoint file with text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "This is a test slide"
    prs.save(files['valid_pptx_file'])

    # Create a valid PDF file with text
    pdf_document = fitz.open()
    page = pdf_document.new_page()
    page.insert_text((100, 100), "This is a test page")
    pdf_document.save(files['valid_pdf_file'])

    # Create synthetic invalid data (empty files)
    files['empty_word_file'] = 'empty_test.docx'
    files['empty_excel_file'] = 'empty_test.xlsx'
    files['empty_pptx_file'] = 'empty_test.pptx'
    files['empty_pdf_file'] = 'empty_test.pdf'
    
    open(files['empty_word_file'], 'w').close()
    open(files['empty_excel_file'], 'w').close()
    open(files['empty_pptx_file'], 'w').close()
    open(files['empty_pdf_file'], 'w').close()
    
    # Create synthetic invalid data (corrupted files)
    files['corrupted_word_file'] = 'corrupted_test.docx'
    files['corrupted_excel_file'] = 'corrupted_test.xlsx'
    files['corrupted_pptx_file'] = 'corrupted_test.pptx'
    files['corrupted_pdf_file'] = 'corrupted_test.pdf'
    
    with open(files['corrupted_word_file'], 'wb') as f:
        f.write(b'\x00\x00\x00\x00')
    with open(files['corrupted_excel_file'], 'wb') as f:
        f.write(b'\x00\x00\x00\x00')
    with open(files['corrupted_pptx_file'], 'wb') as f:
        f.write(b'\x00\x00\x00\x00')
    with open(files['corrupted_pdf_file'], 'wb') as f:
        f.write(b'\x00\x00\x00\x00')
    
    yield files
    
    # Cleanup
    for file in files.values():
        if os.path.exists(file):
            os.remove(file)

# Test for DataProcessingManager

@patch('data_processing.data_processing_manager.DatabaseManager')
@patch('data_processing.data_processing_manager.HashChecker')
@patch('data_processing.data_processing_manager.extract_info')
@patch('data_processing.data_processing_manager.AppInitialization')
def test_process_data(mock_app_init, mock_extract_info, mock_hash_checker, mock_db_manager, setup_files):
    mock_hash_checker.check_hash_exists.return_value = False
    mock_extract_info.return_value = {
        'timeframes': [],
        'locations': [],
        'subjects': [],
        'topics': [],
        'keywords': [],
        'MGRS': [],
        'highest_classification': 'none_found',
        'caveats': 'none_found',
        'full_text': "This is a test document."
    }
    mock_db_manager_instance = MagicMock()
    mock_db_manager.return_value = mock_db_manager_instance

    manager = DataProcessingManager()
    file_path_dict = {
        'docx': [setup_files['valid_word_file']],
        'xlsx': [setup_files['valid_excel_file']],
        'pptx': [setup_files['valid_pptx_file']],
        'pdf': [setup_files['valid_pdf_file']]
    }

    manager.process_data(file_path_dict)

    assert mock_extract_info.called
    assert mock_db_manager_instance.save_data.called



@patch('data_processing.data_processing_manager.AppInitialization')
@patch('data_processing.data_processing_manager.HashChecker')
@patch('data_processing.data_processing_manager.progress_window.ProgressWindow', new=MagicMock)
def test_process_data_with_empty_files(mock_hash_checker, mock_app_init, setup_files):
    mock_hash_checker.check_hash_exists.return_value = False
    mock_app_init.return_value = MagicMock()

    manager = DataProcessingManager()
    file_path_dict = {
        'docx': [setup_files['empty_word_file']],
        'xlsx': [setup_files['empty_excel_file']],
        'pptx': [setup_files['empty_pptx_file']],
        'pdf': [setup_files['empty_pdf_file']]
    }

    with patch('data_extraction.text_extractors.extract_text_from_word', side_effect=PackageNotFoundError):
        with patch('data_extraction.text_extractors.extract_text_from_excel', side_effect=BadZipFile):
            with patch('data_extraction.text_extractors.extract_text_from_pptx', side_effect=PptxPackageNotFoundError):
                with patch('data_extraction.text_extractors.extract_text_from_pdf', side_effect=pdfminer.pdfparser.PDFSyntaxError):
                    for file_type, file_list in file_path_dict.items():
                        for file in file_list:
                            if file_type == 'xlsx':
                                with pytest.raises(BadZipFile):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            elif file_type == 'pdf':
                                with pytest.raises(pdfminer.pdfparser.PDFSyntaxError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            elif file_type == 'pptx':
                                with pytest.raises(PptxPackageNotFoundError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            else:
                                with pytest.raises(PackageNotFoundError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})

@patch('data_processing.data_processing_manager.AppInitialization')
@patch('data_processing.data_processing_manager.HashChecker')
@patch('data_processing.data_processing_manager.progress_window.ProgressWindow', new=MagicMock)
def test_process_data_with_corrupted_files(mock_hash_checker, mock_app_init, setup_files):
    mock_hash_checker.check_hash_exists.return_value = False
    mock_app_init.return_value = MagicMock()

    manager = DataProcessingManager()
    file_path_dict = {
        'docx': [setup_files['corrupted_word_file']],
        'xlsx': [setup_files['corrupted_excel_file']],
        'pptx': [setup_files['corrupted_pptx_file']],
        'pdf': [setup_files['corrupted_pdf_file']]
    }

    with patch('data_extraction.text_extractors.extract_text_from_word', side_effect=PackageNotFoundError):
        with patch('data_extraction.text_extractors.extract_text_from_excel', side_effect=BadZipFile):
            with patch('data_extraction.text_extractors.extract_text_from_pptx', side_effect=PptxPackageNotFoundError):
                with patch('data_extraction.text_extractors.extract_text_from_pdf', side_effect=pdfminer.pdfparser.PDFSyntaxError):
                    for file_type, file_list in file_path_dict.items():
                        for file in file_list:
                            if file_type == 'xlsx':
                                with pytest.raises(BadZipFile):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            elif file_type == 'pdf':
                                with pytest.raises(pdfminer.pdfparser.PDFSyntaxError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            elif file_type == 'pptx':
                                with pytest.raises(PptxPackageNotFoundError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
                            else:
                                with pytest.raises(PackageNotFoundError):
                                    print(f"Testing {file_type} file: {file}")
                                    manager.process_data({file_type: [file]})
# Test for info_processing

@patch('data_processing.info_processing.AppInitialization')
def test_extract_info_valid_text(mock_app_init):
    mock_app_init.nlp = MagicMock()
    mock_app_init.nlp.return_value = MagicMock()
    mock_app_init.nlp.return_value.ents = [
        MagicMock(label_='GPE', text='New York'),
        MagicMock(label_='DATE', text='January 1, 2020')
    ]

    text = "This is a sample text with a location New York and a date January 1, 2020."
    info = extract_info(text)
    assert info['highest_classification'] is not None
    assert 'New York' in info['locations']
    assert 'January 1, 2020' in info['timeframes']

@patch('data_processing.info_processing.AppInitialization')
@patch('data_processing.info_processing.extract_topics')
def test_extract_info_empty_text(mock_extract_topics, mock_app_init):
    mock_app_init.nlp = MagicMock()
    mock_app_init.nlp.return_value = MagicMock()
    mock_app_init.nlp.return_value.ents = []

    # Mock extract_topics to return an empty list for empty text
    mock_extract_topics.return_value = []

    text = ""
    info = extract_info(text)

    assert info['topics'] == [], f"Expected empty list for topics, got {info['topics']}"
    assert info['highest_classification'] == 'none_found', f"Expected 'none_found' for classifications, got {info['classifications']}"
    assert info['caveats'] == 'none_found', f"Expected 'none_found' for caveats, got {info['caveats']}"


def test_extract_classification_and_caveats():
    text = "This document is classified (S//FOUO) with caveats (C//NOFORN)."
    classification, caveats = extract_classification_and_caveats(text)
    assert classification == "S", f"Expected 'S', got {classification}"
    assert "FOUO" in caveats, f"Expected 'FOUO' in caveats, got {caveats}"
    assert "NOFORN" in caveats, f"Expected 'NOFORN' in caveats, got {caveats}"

@patch('utilities.keyword_loading.Keywords.load_latest_keywords')
def test_extract_keywords(mock_load_keywords):
    keywords_list = [
        "Assess",
        "Assessed",
        "Capacity",
        "Capability assessment",
        "Refugee",
        "Refugees",
        "CIMIC",
        "Processes",
        "Developing",
        "Developed",
        "Relationships",
        "Build Relationships",
        "Livelihood",
        "Livelihoods",
        "Collaborate",
        "Build Partner Capacity",
        "Civil network",
        "Hybrid networks",
        "Persistent engagement",
        "Meeting working group",
        "Shipment",
        "Vulnerable population",
        "Influence",
        "Narrative",
        "Information space",
        "Civil reconnaissance",
        "Civil engagement",
        "Non-governmental organization",
        "NGO",
        "United States agency for international development",
        "USAID",
        "Overseas humanitarian, disaster, and civic aid",
        "OHDACA",
        "OHDACA project",
        "Project proposal",
        "Austere medicine",
        "Agriculture",
        "Training",
        "Drug",
        "Drug smuggling",
        "Illicit activity",
        "Cross border",
        "Infrastructures",
        "Strengthen",
        "Partner",
        "School",
        "Students",
        "Gender based violence",
        "Civil Society",
        "Civil Society organization",
        "Civil Society organizations",
        "Crisis",
        "Crime",
        "Instability",
        "Stability",
        "Economic",
        "Initiative",
        "Initiatives",
        "Societal Norms",
        "Civil grievances",
        "Domestic Violence",
        "Educational System Shortfalls",
        "Deteriorating",
        "Ambassador",
        "Logistics hub",
        "Human infrastructure",
        "Unemployment",
        "Assistance with funding"
    ]

    mock_load_keywords.return_value = keywords_list
    Keywords.keywords_list = keywords_list  # Ensure the keywords list is set correctly
#