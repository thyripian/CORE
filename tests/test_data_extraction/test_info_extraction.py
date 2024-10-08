import os
from io import BytesIO

import fitz  # PyMuPDF
import pandas as pd
import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.drawing.image import Image as xlImage
from PIL import Image
from pptx import Presentation

from data_extraction.image_extractors import (extract_images_from_excel,
                                              extract_images_from_pdf,
                                              extract_images_from_pptx,
                                              extract_images_from_word)
from data_extraction.text_extractors import (extract_text_from_excel,
                                             extract_text_from_pdf,
                                             extract_text_from_pptx,
                                             extract_text_from_word)


@pytest.fixture(scope="module")
def setup_files():
    files = {}

    # Create synthetic valid data
    files["valid_word_file"] = "valid_test.docx"
    files["valid_excel_file"] = "valid_test.xlsx"
    files["valid_pptx_file"] = "valid_test.pptx"
    files["valid_pdf_file"] = "valid_test.pdf"

    # Create a valid Word file with an image
    doc = Document()
    doc.add_paragraph("This is a test document.")
    image_stream = BytesIO()
    image = Image.new("RGB", (100, 100), color="red")
    image.save(image_stream, format="PNG")
    image_stream.seek(0)
    doc.add_picture(image_stream, width=1000000)
    doc.save(files["valid_word_file"])

    # Create a valid Excel file with an image
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "This is a test"
    image_stream = BytesIO()
    image = Image.new("RGB", (100, 100), color="red")
    image.save(image_stream, format="PNG")
    image_stream.seek(0)
    img = xlImage(image_stream)
    ws.add_image(img, "B2")
    wb.save(files["valid_excel_file"])

    # Create a valid PowerPoint file with an image and text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "This is a test slide"
    image_stream = BytesIO()
    image.save(image_stream, format="PNG")
    image_stream.seek(0)
    slide.shapes.add_picture(image_stream, 0, 0, width=1000000, height=1000000)
    prs.save(files["valid_pptx_file"])

    # Create a valid PDF file with an image and text
    pdf_document = fitz.open()
    page = pdf_document.new_page()
    rect = fitz.Rect(100, 100, 200, 200)
    image_stream = BytesIO()
    image.save(image_stream, format="PNG")
    image_stream.seek(0)
    page.insert_image(rect, stream=image_stream.getvalue())
    page.insert_text((100, 50), "This is a test page")
    pdf_document.save(files["valid_pdf_file"])

    # Create synthetic invalid data (empty files)
    files["empty_word_file"] = "empty_test.docx"
    files["empty_excel_file"] = "empty_test.xlsx"
    files["empty_pptx_file"] = "empty_test.pptx"
    files["empty_pdf_file"] = "empty_test.pdf"

    open(files["empty_word_file"], "w").close()
    open(files["empty_excel_file"], "w").close()
    open(files["empty_pptx_file"], "w").close()
    open(files["empty_pdf_file"], "w").close()

    # Create synthetic invalid data (corrupted files)
    files["corrupted_word_file"] = "corrupted_test.docx"
    files["corrupted_excel_file"] = "corrupted_test.xlsx"
    files["corrupted_pptx_file"] = "corrupted_test.pptx"
    files["corrupted_pdf_file"] = "corrupted_test.pdf"

    with open(files["corrupted_word_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_excel_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_pptx_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")
    with open(files["corrupted_pdf_file"], "wb") as f:
        f.write(b"\x00\x00\x00\x00")

    yield files

    # Cleanup
    for file in files.values():
        if os.path.exists(file):
            os.remove(file)


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("valid_word_file", extract_images_from_word),
        ("valid_excel_file", extract_images_from_excel),
        ("valid_pptx_file", extract_images_from_pptx),
        ("valid_pdf_file", extract_images_from_pdf),
    ],
)
def test_extract_images_valid(setup_files, file_key, extract_function):
    images = extract_function(setup_files[file_key])
    assert len(images) > 0


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("empty_word_file", extract_images_from_word),
        ("empty_excel_file", extract_images_from_excel),
        ("empty_pptx_file", extract_images_from_pptx),
        ("empty_pdf_file", extract_images_from_pdf),
    ],
)
def test_extract_images_empty(setup_files, file_key, extract_function):
    with pytest.raises(Exception):
        extract_function(setup_files[file_key])


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("corrupted_word_file", extract_images_from_word),
        ("corrupted_excel_file", extract_images_from_excel),
        ("corrupted_pptx_file", extract_images_from_pptx),
        ("corrupted_pdf_file", extract_images_from_pdf),
    ],
)
def test_extract_images_corrupted(setup_files, file_key, extract_function):
    with pytest.raises(Exception):
        extract_function(setup_files[file_key])


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("valid_word_file", extract_text_from_word),
        ("valid_excel_file", extract_text_from_excel),
        ("valid_pptx_file", extract_text_from_pptx),
        ("valid_pdf_file", extract_text_from_pdf),
    ],
)
def test_extract_text_valid(setup_files, file_key, extract_function):
    text = extract_function(setup_files[file_key])
    assert len(text) > 0


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("empty_word_file", extract_text_from_word),
        ("empty_excel_file", extract_text_from_excel),
        ("empty_pptx_file", extract_text_from_pptx),
        ("empty_pdf_file", extract_text_from_pdf),
    ],
)
def test_extract_text_empty(setup_files, file_key, extract_function):
    with pytest.raises(Exception):
        extract_function(setup_files[file_key])


@pytest.mark.parametrize(
    "file_key,extract_function",
    [
        ("corrupted_word_file", extract_text_from_word),
        ("corrupted_excel_file", extract_text_from_excel),
        ("corrupted_pptx_file", extract_text_from_pptx),
        ("corrupted_pdf_file", extract_text_from_pdf),
    ],
)
def test_extract_text_corrupted(setup_files, file_key, extract_function):
    with pytest.raises(Exception):
        extract_function(setup_files[file_key])
